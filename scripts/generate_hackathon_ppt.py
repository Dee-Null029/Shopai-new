from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parents[1] / "ShopAI_Hackathon_Presentation.pptx"


SLIDES = [
    {
        "title": "ShopAI",
        "subtitle": "AI-Powered Smart Shopping Across Amazon, Flipkart, and Myntra\n\nSearch smarter, compare faster, shop with confidence.",
        "bullets": [],
        "notes": "ShopAI is a unified shopping assistant that brings together product discovery, AI analysis, virtual try-on, and conversational guidance into one experience.",
    },
    {
        "title": "The Problem",
        "subtitle": "Online shopping is fragmented and uncertain",
        "bullets": [
            "Users jump across multiple platforms to compare products",
            "Ratings alone do not reveal real product quality",
            "Fashion purchases have fit uncertainty",
            "Decision-making is slow and overwhelming",
        ],
        "notes": "Today, shoppers manually compare the same product across different sites, read scattered reviews, and still remain unsure about quality, trust, and fit.",
    },
    {
        "title": "Our Solution",
        "subtitle": "One platform. Smarter decisions. Better shopping.",
        "bullets": [
            "Aggregate products from Amazon, Flipkart, and Myntra",
            "Rank products using AI-driven scoring",
            "Analyze customer sentiment from reviews",
            "Enable virtual try-on for fashion products",
            "Provide a fashion chatbot for recommendations",
        ],
        "notes": "Instead of forcing users to open multiple apps and interpret raw information themselves, ShopAI acts as an intelligent layer on top of e-commerce platforms.",
    },
    {
        "title": "Key Features",
        "subtitle": "What makes ShopAI different",
        "bullets": [
            "Multi-platform search with deduplication and relevance filtering",
            "Smart ranking engine based on price, rating, sentiment, and reliability",
            "AI sentiment analysis for review insights",
            "3D virtual try-on with adjustable body parameters",
            "AI fashion assistant for chat-based discovery and recommendations",
        ],
        "notes": "The key differentiator is not just search. It is the combination of comparison, intelligence, fit visualization, and conversational guidance in one workflow.",
    },
    {
        "title": "How It Works",
        "subtitle": "From query to confident purchase",
        "bullets": [
            "User enters a product query",
            "ShopAI scrapes Amazon, Flipkart, and Myntra in parallel",
            "Products are normalized, filtered, and deduplicated",
            "AI analyzes review sentiment",
            "Ranking engine computes a final score",
            "User explores results, try-on, and chatbot recommendations",
        ],
        "notes": "This pipeline allows us to turn raw marketplace data into a decision-support experience rather than a basic search page.",
    },
    {
        "title": "Technical Architecture",
        "subtitle": "Built for intelligence, speed, and modularity",
        "bullets": [
            "Frontend: React, Vite, TailwindCSS, Three.js",
            "Backend: Node.js, Express, Socket.io",
            "Data Layer: MongoDB and Redis",
            "Scraping: Puppeteer-based marketplace scrapers",
            "AI Layer: OpenAI-powered sentiment, chat, and body estimation",
            "Deployment: Docker, Nginx, GitHub Actions CI",
        ],
        "notes": "The architecture is modular, which means each major capability such as search, ranking, sentiment, chat, and try-on can evolve independently while sharing a common API layer.",
    },
    {
        "title": "Why It Matters",
        "subtitle": "User value and hackathon impact",
        "bullets": [
            "Saves time in cross-platform comparison",
            "Improves buying confidence through AI insights",
            "Reduces uncertainty in apparel shopping",
            "Creates a richer shopping experience than price sorting alone",
            "Opens future monetization via affiliate links and analytics",
        ],
        "notes": "ShopAI is not just a technical demo. It addresses a real shopping pain point while also creating room for business value through affiliate commerce and personalization.",
    },
    {
        "title": "Demo Flow",
        "subtitle": "Live demo plan",
        "bullets": [
            "Search for a product across all supported platforms",
            "Show ranked results and explain the scoring logic",
            "Open a product and highlight sentiment analysis",
            "Switch to the try-on experience and adjust body parameters",
            "Show the AI chatbot recommending products or styles",
        ],
        "notes": "The demo should feel like one continuous shopping journey: discovery, evaluation, personalization, and recommendation.",
    },
    {
        "title": "Current Progress",
        "subtitle": "What is already working",
        "bullets": [
            "Multi-platform search and product normalization",
            "Smart ranking engine with customizable weights",
            "Sentiment analysis backend with fallback logic",
            "3D try-on viewer and body scanning support",
            "Real-time chat backend with product search integration",
            "Authentication, rate limiting, and secure affiliate tracking",
        ],
        "notes": "For the hackathon, the strongest point is that this is not just an idea. The core platform capabilities are already implemented and connected.",
    },
    {
        "title": "Roadmap",
        "subtitle": "What comes next",
        "bullets": [
            "Complete polished chat UI frontend flow",
            "Add real garment 3D models for richer try-on",
            "Expand admin analytics and user insights",
            "Add password reset, email verification, and stronger testing",
            "Improve production readiness and scale scraping reliability",
        ],
        "notes": "We are being honest about what is complete and what is next. That makes the project more credible during judging.",
    },
    {
        "title": "Closing",
        "subtitle": "ShopAI brings intelligence to online shopping",
        "bullets": [
            "Compare products across platforms",
            "Understand reviews instantly",
            "Personalize fashion decisions",
            "Shop with more confidence",
            "ShopAI transforms online shopping from information overload into intelligent decision-making.",
        ],
        "notes": "Close with the idea that ShopAI is not just a search tool. It is an AI-assisted commerce experience.",
    },
]


def apply_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(246, 242, 235)

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.65),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(23, 37, 84)
    shape.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(7.15),
        Inches(13.333),
        Inches(0.35),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(234, 88, 12)
    accent.line.fill.background()


def add_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.4), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(15, 23, 42)

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.6), Inches(0.9))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(71, 85, 105)


def add_bullets(slide, bullets):
    if not bullets:
        return

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(4.2))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.clear()

    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        p.bullet = True
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(30, 41, 59)


def add_note(slide, note_text):
    notes_text_frame = slide.notes_slide.notes_text_frame
    notes_text_frame.text = note_text


def build_deck():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        apply_background(slide)
        add_title(slide, slide_data["title"], slide_data["subtitle"])
        add_bullets(slide, slide_data["bullets"])
        add_note(slide, slide_data["notes"])

    presentation.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")


if __name__ == "__main__":
    build_deck()